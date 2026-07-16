from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, Any

from ingestion.models import RawDocument
from ingestion.loaders.base import BaseLoader


class PPTXLoader(BaseLoader):
    """
    Loader for Microsoft PowerPoint (.pptx) files.
    Extracts text natively from titles, shapes, tables, and speaker notes.
    Preserves exact slide boundaries, tracks rich media stats (images, charts, 
    SmartArt), and pulls native presentation metadata.
    """

    @property
    def name(self) -> str:
        return "pptx"

    def load(self) -> RawDocument:
      
        path = self._validate_file([".pptx"])
        metadata = self._build_base_metadata(path)

        slides_text = []

       
        stats = {
            "total_slides": 0,
            "text_slide_count": 0,
            "empty_slide_count": 0,
            "shape_count": 0,
            "table_count": 0,
            "image_count": 0,
            "chart_count": 0,
            "smartart_count": 0
        }

      
        try:
            presentation = Presentation(path)

            if presentation.core_properties:
                self._extract_pptx_metadata(presentation.core_properties, metadata)

            stats["total_slides"] = len(presentation.slides)

            for i, slide in enumerate(presentation.slides):
                slide_content = []

                
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        slide_content.append(f"# {title_text}")

              
                for shape in slide.shapes:
                    stats["shape_count"] += 1
                    
                    
                    if shape == slide.shapes.title:
                        continue

                 
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        stats["image_count"] += 1
                    elif shape.has_chart:
                        stats["chart_count"] += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.SMART_ART:
                        stats["smartart_count"] += 1

                    
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            slide_content.append(text)
                    
                    
                    if shape.has_table:
                        stats["table_count"] += 1
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_data:
                                slide_content.append(" | ".join(row_data))

             
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_content.append(f"\n[NOTES]\n{notes}\n[/NOTES]")

                joined_slide_text = "\n\n".join(slide_content).strip()

                if joined_slide_text:
                    slides_text.append(joined_slide_text)
                    stats["text_slide_count"] += 1
                else:
                   
                    slides_text.append("")
                    stats["empty_slide_count"] += 1
                    self.log_warning(f"Slide {i + 1} yielded no text (possible image-only or empty slide).")

        except Exception as e:
            raise RuntimeError(f"Fatal error reading PPTX file: {str(e)}")

       
        content = "\n\n--- SLIDE BREAK ---\n\n".join(slides_text).strip()

        if not content:
            self.log_warning("No textual content extracted from PPTX document. It may contain only images.")

        
        metadata.update(stats)
        metadata.update(self._calculate_text_statistics(content))

       
        if stats["total_slides"] > 0:
            metadata["average_characters_per_slide"] = metadata["character_count"] // stats["total_slides"]
        else:
            metadata["average_characters_per_slide"] = 0

        return RawDocument(
            content=content,
            source=self.name,
            metadata=metadata,
            slides=slides_text,
            warnings=self.warnings
        )

    def _extract_pptx_metadata(self, core_props: Any, metadata: Dict[str, Any]):
        """Helper to safely extract native PPTX metadata fields."""
        try:
            if core_props.author:
                metadata["author"] = core_props.author.strip()
            if core_props.title:
                metadata["title"] = core_props.title.strip()
            if core_props.subject:
                metadata["subject"] = core_props.subject.strip()
            if core_props.category:
                metadata["category"] = core_props.category.strip()
            if core_props.comments:
                metadata["comments"] = core_props.comments.strip()
            if core_props.created:
                metadata["created_at"] = core_props.created.isoformat()
        except Exception as e:
            self.log_warning(f"Could not extract native PPTX core properties: {str(e)}")